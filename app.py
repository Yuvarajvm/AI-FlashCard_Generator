import os
import fitz  # PyMuPDF
import pandas as pd
from flask import Flask, render_template, request, send_file, session
import google.generativeai as genai
from docx import Document
from pptx import Presentation
from fpdf import FPDF
from werkzeug.utils import secure_filename

app = Flask(__name__)
app.config["UPLOAD_FOLDER"] = "uploads"
app.config["EXPORT_FOLDER"] = "exports"

# Add secret key for session management - REQUIRED for sessions to work
app.secret_key = os.getenv("FLASK_SECRET_KEY", "default_secret_key_for_dev")

# Configure Gemini API
GEMINI_API_KEY = os.getenv("GEMINI_API_KEY")
genai.configure(api_key=GEMINI_API_KEY)

ALLOWED_EXTENSIONS = {"txt", "pdf", "docx", "pptx"}

def allowed_file(filename):
    return "." in filename and filename.rsplit(".", 1)[1].lower() in ALLOWED_EXTENSIONS

def extract_text_from_file(filepath):
    ext = filepath.rsplit(".", 1)[1].lower()
    text = ""

    if ext == "pdf":
        pdf_doc = fitz.open(filepath)
        for page in pdf_doc:
            text += page.get_text()

    elif ext == "docx":
        doc = Document(filepath)
        for para in doc.paragraphs:
            text += para.text + "\n"

    elif ext == "pptx":
        prs = Presentation(filepath)
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text += shape.text + "\n"

    elif ext == "txt":
        with open(filepath, "r", encoding="utf-8") as f:
            text = f.read()

    return text.strip()

def generate_flashcards(input_text, count):
    model = genai.GenerativeModel("gemini-1.5-flash")
    prompt = f"""
    Generate {count} flashcards from the following content:

    {input_text}

    Return each flashcard in the format:
    Q: question
    A: answer
    """

    response = model.generate_content(prompt)
    output = response.text.strip()

    cards = []
    for pair in output.split("Q:")[1:]:
        if "A:" in pair:
            q, a = pair.split("A:", 1)
            cards.append((q.strip(), a.strip()))
    return cards

def export_to_pdf(cards, filename):
    pdf = FPDF()
    pdf.add_page()

    # Register and use the Unicode font (requires DejaVuSans.ttf in fonts/)
    pdf.add_font('DejaVu', '', 'fonts/DejaVuSans.ttf', uni=True)
    pdf.set_font('DejaVu', '', 12)

    for i, (q, a) in enumerate(cards, 1):
        pdf.multi_cell(0, 8, f"Q{i}: {q}")
        pdf.multi_cell(0, 8, f"A{i}: {a}")
        pdf.ln()
    pdf_path = os.path.join(app.config["EXPORT_FOLDER"], filename)
    pdf.output(pdf_path)
    return pdf_path


def export_to_docx(cards, filename):
    doc = Document()
    for i, (q, a) in enumerate(cards, 1):
        doc.add_paragraph(f"Q{i}: {q}")
        doc.add_paragraph(f"A{i}: {a}")
        doc.add_paragraph("")
    doc_path = os.path.join(app.config["EXPORT_FOLDER"], filename)
    doc.save(doc_path)
    return doc_path

def export_to_csv(cards, filename):
    df = pd.DataFrame(cards, columns=["Question", "Answer"])
    csv_path = os.path.join(app.config["EXPORT_FOLDER"], filename)
    df.to_csv(csv_path, index=False)
    return csv_path

@app.route("/", methods=["GET", "POST"])
def index():
    if request.method == "POST":
        count = int(request.form["count"])
        text_input = request.form.get("text", "").strip()
        file = request.files.get("file")
        extracted_text = ""

        if file and allowed_file(file.filename):
            filename = secure_filename(file.filename)
            filepath = os.path.join(app.config["UPLOAD_FOLDER"], filename)
            file.save(filepath)
            extracted_text = extract_text_from_file(filepath)

        if text_input:
            extracted_text += "\n" + text_input

        if not extracted_text.strip():
            return "No content provided", 400

        flashcards = generate_flashcards(extracted_text, count)
        session["flashcards"] = flashcards  # Fixed: Use Flask session instead of request
        return render_template("result.html", flashcards=flashcards)

    return render_template("index.html")

@app.route("/export/<filetype>")  # Fixed: Correct route parameter syntax
def export(filetype):
    cards = session.get("flashcards", None)  # Fixed: Get from session instead of request
    if not cards:
        return "No cards to export", 400

    if filetype == "pdf":
        path = export_to_pdf(cards, "flashcards.pdf")
    elif filetype == "docx":
        path = export_to_docx(cards, "flashcards.docx")
    elif filetype == "csv":
        path = export_to_csv(cards, "flashcards.csv")
    else:
        return "Invalid format", 400

    return send_file(path, as_attachment=True)

if __name__ == "__main__":
    os.makedirs(app.config["UPLOAD_FOLDER"], exist_ok=True)
    os.makedirs(app.config["EXPORT_FOLDER"], exist_ok=True)
    app.run(host="0.0.0.0", port=int(os.environ.get("PORT", 5000)))
